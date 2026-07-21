from __future__ import annotations

import tempfile
from pathlib import Path

import fitz
import openpyxl
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from reportlab.pdfgen import canvas


def verify_docx(path: Path) -> None:
    document = Document()
    document.add_heading("Office/PDF environment smoke test", level=1)
    document.add_paragraph("DOCX creation works.")
    document.save(path)

    loaded = Document(path)
    text = "\n".join(paragraph.text for paragraph in loaded.paragraphs)
    assert "DOCX creation works." in text


def verify_xlsx(path: Path) -> None:
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Smoke"
    sheet["A1"] = "Library"
    sheet["B1"] = "Status"
    sheet["A2"] = "openpyxl"
    sheet["B2"] = "OK"
    workbook.save(path)

    loaded = openpyxl.load_workbook(path)
    assert loaded["Smoke"]["B2"].value == "OK"

    frame = pd.DataFrame([{"Library": "pandas", "Status": "OK"}])
    pandas_path = path.with_name("pandas-smoke.xlsx")
    frame.to_excel(pandas_path, index=False)
    assert pandas_path.exists()


def verify_pptx(path: Path) -> None:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "PowerPoint smoke test"
    textbox = slide.shapes.add_textbox(914400, 1828800, 7315200, 914400)
    textbox.text_frame.text = "PPTX creation works."
    deck.save(path)

    loaded = Presentation(path)
    assert len(loaded.slides) == 1


def verify_pdf(path: Path) -> None:
    pdf = canvas.Canvas(str(path))
    pdf.drawString(72, 720, "PDF creation and extraction works.")
    pdf.save()

    reader = PdfReader(path)
    assert len(reader.pages) == 1
    assert "PDF creation" in (reader.pages[0].extract_text() or "")

    with pdfplumber.open(path) as opened:
        text = opened.pages[0].extract_text() or ""
        assert "PDF creation" in text

    with fitz.open(path) as opened:
        assert opened.page_count == 1


def main() -> None:
    with tempfile.TemporaryDirectory(prefix="office-pdf-smoke-") as tmp:
        root = Path(tmp)
        verify_docx(root / "smoke.docx")
        verify_xlsx(root / "smoke.xlsx")
        verify_pptx(root / "smoke.pptx")
        verify_pdf(root / "smoke.pdf")

    print("Office/PDF environment OK")
    print("Libraries: python-docx, openpyxl, pandas, python-pptx, pypdf, pdfplumber, pymupdf, reportlab, pillow")


if __name__ == "__main__":
    main()
